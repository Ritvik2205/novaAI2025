from __future__ import annotations

import io
from pathlib import Path
from typing import Iterable

import pandas as pd
from bs4 import BeautifulSoup
from pdfminer.high_level import extract_text
from pptx import Presentation
from pypdf import PdfReader
from docx import Document as DocxDocument
import pytesseract
from PIL import Image

from app.utils.html import extract_main_content


class ParsedDoc:
    def __init__(self, text: str, meta: dict[str, str]) -> None:
        self.text = text
        self.meta = meta


def parse_pdf(path: Path) -> ParsedDoc:
    reader = PdfReader(str(path))
    text_parts = [page.extract_text() or "" for page in reader.pages]
    text = "\n".join(text_parts)
    if not text.strip():
        text = extract_text(str(path))
    if not text.strip():
        image = Image.open(path)
        text = pytesseract.image_to_string(image)
    return ParsedDoc(text=text, meta={"doc_type": "pdf"})


def parse_docx(path: Path) -> ParsedDoc:
    doc = DocxDocument(str(path))
    text = "\n".join(para.text for para in doc.paragraphs)
    return ParsedDoc(text=text, meta={"doc_type": "docx"})


def parse_pptx(path: Path) -> ParsedDoc:
    pres = Presentation(str(path))
    text = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return ParsedDoc(text="\n".join(text), meta={"doc_type": "pptx"})


def parse_csv(path: Path) -> ParsedDoc:
    df = pd.read_csv(path)
    text = df.to_markdown(index=False)
    return ParsedDoc(text=text, meta={"doc_type": "csv"})


def parse_html(content: str) -> ParsedDoc:
    text, meta = extract_main_content(content)
    meta.setdefault("doc_type", "html")
    return ParsedDoc(text=text, meta=meta)


def parse_any(path: Path | str, mime: str | None = None) -> ParsedDoc:
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix in {".docx"}:
        return parse_docx(path)
    if suffix in {".pptx"}:
        return parse_pptx(path)
    if suffix in {".csv"}:
        return parse_csv(path)
    if suffix in {".html", ".htm"}:
        return parse_html(path.read_text())
    raise ValueError(f"Unsupported file {suffix}")
